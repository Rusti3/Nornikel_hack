from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import fitz
import pandas as pd
from docx import Document
from pptx import Presentation

from src.mekg.config import MEKGConfig
from src.mekg.extractor import MEKGExtractor
from src.mekg.models import (
    ChunkExtraction,
    ConditionCandidate,
    ElementKind,
    EntityCandidate,
    EvidenceRef,
    MeasurementCandidate,
    SourceElement,
)
from src.mekg.service import MEKGService
from src.mekg.parsers import DocumentParser
from src.mekg.units import UnitNormalizer


class DisabledVision:
    enabled = False


def config(tmp_path: Path) -> MEKGConfig:
    ontology = Path(__file__).resolve().parents[1] / "ontology"
    return MEKGConfig(
        ontology_dir=ontology,
        artifacts_dir=tmp_path / "artifacts",
        llm_model="test",
        vision_model="test",
        yandex_api_key="",
        yandex_folder_id="",
        yandex_base_url="https://example.invalid/v1",
        yandex_ocr_url="https://example.invalid/ocr",
        data_logging=False,
        max_concurrency=1,
        chunk_chars=2000,
        chunk_overlap=100,
        ocr_min_chars=100,
    )


def test_unit_normalization_ru_and_ranges():
    units = UnitNormalizer()
    assert units.normalize("мг/дм³").symbol == "mg/L"
    assert units.normalize("А/м²").dimension == "current_density"
    assert units.normalize("°C").symbol == "degC"
    assert not units.normalize("totally-unknown-unit").valid


def test_graph_source_query_tokens_keep_technical_anchors_and_drop_question_words():
    tokens = MEKGService._graph_query_tokens(
        "Какие извлечения Co 91, Ni 85 и Mn 60 заявлены для процесса Cuprion?"
    )
    assert "cuprion" in tokens
    assert "91" in tokens
    assert "какие" not in tokens
    assert "процесса" not in tokens


def test_graph_source_query_tokens_add_bilingual_metallurgy_terms():
    tokens = MEKGService._graph_query_tokens(
        "Как изменилась мировая сортность медной руды по регионам в 2012 году?"
    )
    assert {"copper", "ore", "grade", "region"}.issubset(tokens)
    assert "2012" in tokens


def test_extract_validation_separates_invalid_evidence(tmp_path):
    extractor = MEKGExtractor.__new__(MEKGExtractor)
    extractor.units = UnitNormalizer()
    element = SourceElement(
        id="chunk_1",
        kind=ElementKind.TEXT,
        text="At 25 °C nickel recovery reached 94 percent.",
    )
    evidence = EvidenceRef(element_id="chunk_1", quote="nickel recovery reached 94 percent")
    invalid_evidence = EvidenceRef(element_id="other", quote="made up")
    extraction = ChunkExtraction(
        entities=[EntityCandidate(local_id="e1", entity_type="Material", canonical_name="nickel", evidence=evidence)],
        conditions=[ConditionCandidate(local_id="c1", name="temperature", value=25, unit_original="°C", evidence=evidence)],
        measurements=[
            MeasurementCandidate(local_id="m1", name="recovery", property_name="recovery", value=94, unit_original="%", evidence=evidence),
            MeasurementCandidate(local_id="m2", name="pressure", property_name="pressure", value=10, unit_original="MPa", evidence=invalid_evidence),
        ],
    )
    valid, rejected = extractor.validate(extraction, element)
    assert len(valid.entities) == 1
    assert len(valid.conditions) == 1
    assert [item.local_id for item in valid.measurements] == ["m1"]
    assert len(rejected) == 1


def test_entity_type_preserves_controlled_pascal_case():
    evidence = EvidenceRef(element_id="chunk_1", quote="SiO2")
    entity = EntityCandidate(
        local_id="e1",
        entity_type="Chemicalcompound",
        canonical_name="SiO2",
        evidence=evidence,
    )
    assert entity.entity_type == "ChemicalCompound"


def test_yandex_schema_variants_are_normalized_without_losing_response():
    evidence = {"element_id": "chunk_1", "quote": "Pd was approximately 12 percent"}
    entity = EntityCandidate.model_validate(
        {
            "local_id": "e1",
            "entity_type": "Element",
            "name_ru": "палладий",
            "name_en": "palladium",
            "evidence": evidence,
        }
    )
    measurement = MeasurementCandidate.model_validate(
        {
            "local_id": "m1",
            "name": "Pd content",
            "property_name": "content",
            "value": "12,0",
            "comparator": "≈",
            "unit_original": "%",
            "evidence": evidence,
        }
    )
    assert entity.canonical_name == "палладий"
    assert measurement.comparator == "="
    assert measurement.value == 12.0


def test_tool_text_matching_handles_russian_inflection_and_generic_terms():
    assert MEKGService._text_match_score(
        "никель", ["электроэкстракция никеля из хлоридного электролита"]
    ) >= 0.45
    assert MEKGService._text_match_score(
        "технология очистки шахтных вод",
        ["технология кучного выщелачивания"],
        ignore_generic=True,
    ) < 0.45
    assert MEKGService._text_match_score(
        "мембранная очистка шахтных вод",
        ["мембранная дистилляция для обработки шахтной воды"],
        ignore_generic=True,
    ) >= 0.50
    assert MEKGService._text_match_score(
        "шахтных вод",
        ["технология производит концентрат"],
        ignore_generic=True,
    ) == 0
    assert MEKGService._text_match_score(
        "очистка шахтных вод обратный осмос",
        ["магнетизирующий обжиг в водороде и обратная флотация"],
        ignore_generic=True,
    ) < 0.30
    assert MEKGService._technology_label_from_claim(
        "Обратный осмос удаляет до 80 % нитрата из воды."
    ) == "Обратный осмос"


def test_table_layout_evidence_and_symbolic_numeric_are_handled_per_item():
    extractor = MEKGExtractor.__new__(MEKGExtractor)
    extractor.units = UnitNormalizer()
    element = SourceElement(
        id="slide_1",
        kind=ElementKind.TEXT,
        text="ОХВ ПНТП (ЦЭН-2), %\nS\nSiO2\n7,31-\n15,0\n10,3-\n16,2",
    )
    evidence = EvidenceRef(
        element_id="slide_1",
        quote="ОХВ ПНТП (ЦЭН-2), % S 7,31-15,0",
    )
    extraction = ChunkExtraction(
        entities=[
            EntityCandidate(
                local_id="e1",
                entity_type="Chemicalcompound",
                canonical_name="S",
                evidence=evidence,
            )
        ],
        conditions=[
            ConditionCandidate(
                local_id="c1",
                name="S content",
                value_min="7,31",
                value_max="15,0",
                comparator="range",
                unit_original="%",
                evidence=evidence,
            )
        ],
    )
    valid, rejected = extractor.validate(extraction, element)
    assert valid.entities[0].entity_type == "ChemicalCompound"
    assert valid.conditions[0].value_min == 7.31
    assert valid.conditions[0].evidence.confidence == 0.7
    assert rejected == []

    malformed = ConditionCandidate.model_validate(
        {
            "local_id": "c2",
            "name": "symbolic radius",
            "value": "r + 5R",
            "unit_original": "мм",
            "evidence": evidence.model_dump(),
        }
    )
    invalid, rejected = extractor.validate(ChunkExtraction(conditions=[malformed]), element)
    assert invalid.conditions == []
    assert rejected[0]["reason"] == "non-numeric value"


def test_pdf_docx_pptx_and_xlsx_parsers(tmp_path):
    parser = DocumentParser(config(tmp_path), DisabledVision())

    pdf_path = tmp_path / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Nickel electrowinning experiment at 25 C. " * 8)
    pdf.save(pdf_path)

    docx_path = tmp_path / "sample.docx"
    word = Document()
    word.add_heading("Mine water treatment", 0)
    word.add_paragraph("Sulfate concentration was 200 mg/L. " * 4)
    table = word.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "SO4"
    table.cell(1, 1).text = "200 mg/L"
    word.save(docx_path)

    pptx_path = tmp_path / "sample.pptx"
    slides = Presentation()
    slide = slides.slides.add_slide(slides.slide_layouts[1])
    slide.shapes.title.text = "PGM distribution"
    slide.placeholders[1].text = "Distribution between matte and slag at 1250 C"
    slides.save(pptx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    pd.DataFrame([["Country", "Copper kt"], ["Russia", 100]]).to_excel(xlsx_path, index=False, header=False)

    parsed_pdf = parser.parse(pdf_path)
    parsed_docx = parser.parse(docx_path)
    parsed_pptx = parser.parse(pptx_path)
    parsed_xlsx = parser.parse(xlsx_path)

    assert any(item.kind == ElementKind.TEXT for item in parsed_pdf.elements)
    assert any(item.kind == ElementKind.TABLE_ROW for item in parsed_docx.elements)
    assert any(item.slide_number == 1 for item in parsed_pptx.elements)
    assert any(item.sheet_name == "Sheet1" for item in parsed_xlsx.elements)
    assert len({parsed_pdf.document_id, parsed_docx.document_id, parsed_pptx.document_id, parsed_xlsx.document_id}) == 4


def test_legacy_doc_uses_antiword_text_layer(tmp_path, monkeypatch):
    parser = DocumentParser(config(tmp_path), DisabledVision())
    doc_path = tmp_path / "legacy.doc"
    doc_path.write_bytes(b"fake legacy word content")

    def fake_run(command, **kwargs):
        assert command[0] == "antiword"
        return SimpleNamespace(
            returncode=0,
            stdout="Nickel recovery reached 94 percent at 25 C. " * 4,
            stderr="",
        )

    monkeypatch.setattr("src.mekg.parsers.subprocess.run", fake_run)
    parsed = parser.parse(doc_path, source_locator="folder/legacy.doc")

    assert parsed.file_type == "doc"
    assert any("Nickel recovery reached 94 percent" in item.text for item in parsed.elements)
    assert any("antiword" in warning for warning in parsed.warnings)


def test_same_source_locator_has_stable_document_id(tmp_path):
    parser = DocumentParser(config(tmp_path), DisabledVision())
    path = tmp_path / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Stable identity test text. " * 8)
    pdf.save(path)
    first = parser.parse(path, source_locator="folder/sample.pdf")
    second = parser.parse(path, source_locator="folder/sample.pdf")
    assert first.document_id == second.document_id
    assert first.version_id == second.version_id
