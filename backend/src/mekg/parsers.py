from __future__ import annotations

import hashlib
import json
import logging
import mimetypes
import re
import shutil
import subprocess
import tempfile
import uuid
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Iterable

import fitz
import pandas as pd
from docx import Document as WordDocument
from PIL import Image
from pptx import Presentation

from .config import MEKGConfig
from .models import ElementKind, ParsedDocument, PublicationBoundary, SourceElement
from .vision import YandexVisionClient


NAMESPACE = uuid.UUID("69353598-1dc0-4aeb-98d9-f8ed395da6ad")


def stable_id(prefix: str, value: str) -> str:
    return f"{prefix}_{uuid.uuid5(NAMESPACE, value)}"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def clean_text(value: str) -> str:
    value = value.replace("\x00", " ").replace("\u00ad", "")
    return re.sub(r"[ \t]+", " ", re.sub(r"\r\n?", "\n", value)).strip()


def infer_language(text: str) -> str | None:
    if not text:
        return None
    cyr = len(re.findall(r"[А-Яа-яЁё]", text))
    lat = len(re.findall(r"[A-Za-z]", text))
    if cyr > lat * 1.5:
        return "ru"
    if lat > cyr * 1.5:
        return "en"
    return "mixed"


class DocumentParser:
    SUPPORTED = {".pdf", ".doc", ".docx", ".docm", ".pptx", ".xls", ".xlsx"}

    def __init__(self, config: MEKGConfig | None = None, vision: YandexVisionClient | None = None) -> None:
        self.config = config or MEKGConfig.from_env()
        self.config.artifacts_dir.mkdir(parents=True, exist_ok=True)
        self.vision = vision or YandexVisionClient(self.config)

    def parse(
        self,
        path: str | Path,
        *,
        source_locator: str | None = None,
        category: str | None = None,
        fast: bool = False,
    ) -> ParsedDocument:
        path = Path(path)
        if not path.is_file():
            raise FileNotFoundError(path)
        suffix = path.suffix.lower()
        if suffix not in self.SUPPORTED:
            raise ValueError(f"Unsupported MEKG source format: {suffix}")
        sha = file_sha256(path)
        locator = (source_locator or str(path.resolve())).replace("\\", "/")
        document_id = stable_id("doc", locator.casefold())
        version_id = f"docver_{sha}"
        document = ParsedDocument(
            document_id=document_id,
            version_id=version_id,
            source_locator=locator,
            file_name=path.name,
            file_type=suffix.lstrip("."),
            sha256=sha,
            size_bytes=path.stat().st_size,
            category=category,
        )
        artifact_dir = self.config.artifacts_dir / version_id
        artifact_dir.mkdir(parents=True, exist_ok=True)
        if suffix == ".pdf":
            self._parse_pdf(path, document, artifact_dir, fast=fast)
        elif suffix == ".doc":
            self._parse_doc(path, document)
        elif suffix == ".docm":
            try:
                self._parse_docx(path, document, artifact_dir, fast=fast)
            except Exception as exc:
                if fast:
                    self._parse_openxml_word(path, document)
                    document.warnings.append(f"docm parsed from OpenXML text layer after python-docx {type(exc).__name__}")
                else:
                    converted = self._convert_office(path, artifact_dir, "docx")
                    self._parse_docx(converted, document, artifact_dir, fast=fast)
        elif suffix == ".docx":
            try:
                self._parse_docx(path, document, artifact_dir, fast=fast)
            except Exception as exc:
                if not fast:
                    raise
                self._parse_openxml_word(path, document)
                document.warnings.append(f"docx parsed from OpenXML text layer after python-docx {type(exc).__name__}")
        elif suffix == ".pptx":
            self._parse_pptx(path, document, artifact_dir, fast=fast)
        else:
            self._parse_workbook(path, document)

        textual = "\n".join(element.text for element in document.elements if element.text)
        document.language = infer_language(textual[:100000])
        document.title = self._infer_title(document)
        document.publications = self._initial_publications(document)
        return document

    def _element_id(self, document: ParsedDocument, kind: str, location: str) -> str:
        return stable_id(kind, f"{document.version_id}:{location}")

    def _append_text_chunks(
        self,
        document: ParsedDocument,
        text: str,
        *,
        location: str,
        page_number: int | None = None,
        slide_number: int | None = None,
        sheet_name: str | None = None,
        metadata: dict | None = None,
    ) -> None:
        text = clean_text(text)
        if not text:
            return
        size = self.config.chunk_chars
        overlap = min(self.config.chunk_overlap, size // 3)
        start = 0
        part = 0
        while start < len(text):
            end = min(len(text), start + size)
            if end < len(text):
                boundary = max(text.rfind("\n", start, end), text.rfind(". ", start, end))
                if boundary > start + size // 2:
                    end = boundary + 1
            chunk = text[start:end].strip()
            if chunk:
                document.elements.append(
                    SourceElement(
                        id=self._element_id(document, "chunk", f"{location}:{part}"),
                        kind=ElementKind.TEXT,
                        text=chunk,
                        page_number=page_number,
                        slide_number=slide_number,
                        sheet_name=sheet_name,
                        metadata={**(metadata or {}), "part": part, "char_start": start, "char_end": end},
                    )
                )
            if end >= len(text):
                break
            start = max(start + 1, end - overlap)
            part += 1

    def _parse_pdf(
        self, path: Path, document: ParsedDocument, artifact_dir: Path, *, fast: bool = False
    ) -> None:
        pdf = fitz.open(path)
        image_dir = artifact_dir / "figures"
        if not fast:
            image_dir.mkdir(exist_ok=True)
        seen_images: set[str] = set()
        for index, page in enumerate(pdf):
            page_number = index + 1
            text = clean_text(page.get_text("text") or "")
            if not fast and len(text) < self.config.ocr_min_chars and self.vision.enabled:
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    ocr_bytes = pix.tobytes("png")
                    text = clean_text(self.vision.recognize_text(ocr_bytes, model="page-column-sort"))
                    if text:
                        document.warnings.append(f"page {page_number}: OCR used")
                except Exception as exc:
                    document.warnings.append(f"page {page_number}: OCR failed ({type(exc).__name__})")
                    logging.warning("MEKG OCR failed for %s page %s: %s", path.name, page_number, exc)
            self._append_text_chunks(document, text, location=f"page:{page_number}", page_number=page_number)
            if fast:
                continue
            self._extract_pdf_tables(page, document, page_number)
            for image_index, image in enumerate(page.get_images(full=True)):
                try:
                    extracted = pdf.extract_image(image[0])
                    content = extracted["image"]
                    width = int(extracted.get("width", 0))
                    height = int(extracted.get("height", 0))
                    if width < 200 or height < 120 or len(content) < 8000:
                        continue
                    digest = hashlib.sha256(content).hexdigest()
                    if digest in seen_images:
                        continue
                    seen_images.add(digest)
                    extension = extracted.get("ext", "png")
                    target = image_dir / f"{digest}.{extension}"
                    target.write_bytes(content)
                    document.elements.append(
                        SourceElement(
                            id=self._element_id(document, "figure", f"page:{page_number}:{digest}"),
                            kind=ElementKind.FIGURE,
                            page_number=page_number,
                            image_path=str(target),
                            metadata={"sha256": digest, "width": width, "height": height, "image_index": image_index},
                        )
                    )
                except Exception as exc:
                    document.warnings.append(f"page {page_number}: image extraction failed ({type(exc).__name__})")

    def _extract_pdf_tables(self, page: fitz.Page, document: ParsedDocument, page_number: int) -> None:
        try:
            tables = page.find_tables().tables
        except Exception:
            return
        for table_index, table in enumerate(tables):
            rows = table.extract()
            table_id = self._element_id(document, "table", f"page:{page_number}:{table_index}")
            document.elements.append(
                SourceElement(
                    id=table_id,
                    kind=ElementKind.TABLE,
                    page_number=page_number,
                    bbox=list(table.bbox),
                    text=f"Table with {len(rows)} rows",
                    metadata={"table_index": table_index},
                )
            )
            for row_index, row in enumerate(rows):
                values = [clean_text(str(value or "")) for value in row]
                if not any(values):
                    continue
                document.elements.append(
                    SourceElement(
                        id=self._element_id(document, "row", f"page:{page_number}:{table_index}:{row_index}"),
                        kind=ElementKind.TABLE_ROW,
                        text=" | ".join(values),
                        page_number=page_number,
                        row_number=row_index + 1,
                        metadata={"table_id": table_id, "cells": values},
                    )
                )

    def _parse_doc(self, path: Path, document: ParsedDocument) -> None:
        """Extract legacy binary Word documents through antiword.

        For MEKG ingestion the important contract is the text evidence layer.
        LibreOffice conversion is useful when it works, but this Docker image's
        headless ``soffice`` can be brittle for old ``.doc`` files.  ``antiword``
        gives us a smaller, deterministic text path for those files.
        """
        attempts = [
            ["antiword", "-m", "UTF-8", str(path)],
            ["antiword", str(path)],
        ]
        errors: list[str] = []
        for command in attempts:
            result = subprocess.run(command, capture_output=True, text=True, timeout=180, check=False)
            text = clean_text(result.stdout)
            if result.returncode == 0 and len(text) >= 20:
                self._append_text_chunks(
                    document,
                    text,
                    location="doc-page:1",
                    page_number=1,
                    metadata={"parser": "antiword", "legacy_word": True},
                )
                document.warnings.append("legacy .doc parsed with antiword; tables and images were not extracted")
                return
            errors.append((result.stderr or result.stdout or "").strip()[-500:])
        raise RuntimeError(f"antiword failed for {path.name}: {' | '.join(error for error in errors if error)}")

    def _convert_office(self, path: Path, artifact_dir: Path, target_format: str) -> Path:
        output_dir = artifact_dir / "converted"
        output_dir.mkdir(exist_ok=True)
        command = [
            "soffice",
            "--headless",
            "--convert-to",
            target_format,
            "--outdir",
            str(output_dir),
            str(path),
        ]
        result = subprocess.run(command, capture_output=True, text=True, timeout=180, check=False)
        expected = output_dir / f"{path.stem}.{target_format}"
        if result.returncode != 0 or not expected.exists():
            raise RuntimeError(f"LibreOffice conversion failed for {path.name}: {result.stderr[-500:]}")
        return expected

    def _parse_docx(
        self, path: Path, document: ParsedDocument, artifact_dir: Path, *, fast: bool = False
    ) -> None:
        word = WordDocument(path)
        paragraph_buffer: list[str] = []
        page_number = 1
        for index, paragraph in enumerate(word.paragraphs):
            text = clean_text(paragraph.text)
            if text:
                paragraph_buffer.append(text)
            if "pageBreakBefore" in paragraph._p.xml or len("\n".join(paragraph_buffer)) > self.config.chunk_chars:
                self._append_text_chunks(
                    document,
                    "\n".join(paragraph_buffer),
                    location=f"docx-page:{page_number}",
                    page_number=page_number,
                    metadata={"paragraph_end": index},
                )
                paragraph_buffer = []
                page_number += 1
        if paragraph_buffer:
            self._append_text_chunks(
                document,
                "\n".join(paragraph_buffer),
                location=f"docx-page:{page_number}",
                page_number=page_number,
            )
        for table_index, table in enumerate(word.tables):
            table_id = self._element_id(document, "table", f"docx-table:{table_index}")
            document.elements.append(
                SourceElement(id=table_id, kind=ElementKind.TABLE, text=f"Table {table_index + 1}")
            )
            for row_index, row in enumerate(table.rows):
                cells = [clean_text(cell.text) for cell in row.cells]
                if not any(cells):
                    continue
                document.elements.append(
                    SourceElement(
                        id=self._element_id(document, "row", f"docx-table:{table_index}:{row_index}"),
                        kind=ElementKind.TABLE_ROW,
                        text=" | ".join(cells),
                        row_number=row_index + 1,
                        metadata={"table_id": table_id, "cells": cells},
                    )
                )
        if not fast:
            self._extract_docx_images(path, document, artifact_dir)

    def _extract_docx_images(self, path: Path, document: ParsedDocument, artifact_dir: Path) -> None:
        image_dir = artifact_dir / "figures"
        image_dir.mkdir(exist_ok=True)
        try:
            with zipfile.ZipFile(path) as archive:
                for name in archive.namelist():
                    if not name.startswith("word/media/"):
                        continue
                    content = archive.read(name)
                    if len(content) < 8000:
                        continue
                    digest = hashlib.sha256(content).hexdigest()
                    suffix = Path(name).suffix or ".png"
                    target = image_dir / f"{digest}{suffix}"
                    target.write_bytes(content)
                    try:
                        with Image.open(target) as image:
                            width, height = image.size
                        if width < 200 or height < 120:
                            target.unlink(missing_ok=True)
                            continue
                    except Exception:
                        width = height = 0
                    document.elements.append(
                        SourceElement(
                            id=self._element_id(document, "figure", f"docx:{digest}"),
                            kind=ElementKind.FIGURE,
                            image_path=str(target),
                            metadata={"sha256": digest, "width": width, "height": height},
                        )
                    )
        except (zipfile.BadZipFile, OSError):
            return

    def _parse_openxml_word(self, path: Path, document: ParsedDocument) -> None:
        """Recover text and native table rows without LibreOffice or image extraction."""
        namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        paragraphs: list[str] = []
        with zipfile.ZipFile(path) as archive:
            names = [
                name for name in archive.namelist()
                if name == "word/document.xml" or re.fullmatch(r"word/(?:header|footer)\d+\.xml", name)
            ]
            if "word/document.xml" not in names:
                raise ValueError(f"OpenXML document.xml is missing in {path.name}")
            for name in names:
                root = ET.fromstring(archive.read(name))
                for paragraph in root.findall(".//w:p", namespace):
                    value = clean_text("".join(node.text or "" for node in paragraph.findall(".//w:t", namespace)))
                    if value:
                        paragraphs.append(value)
                for table_index, table in enumerate(root.findall(".//w:tbl", namespace)):
                    table_id = self._element_id(document, "table", f"openxml:{name}:{table_index}")
                    document.elements.append(SourceElement(id=table_id, kind=ElementKind.TABLE, text=f"OpenXML table {table_index + 1}"))
                    for row_index, row in enumerate(table.findall("./w:tr", namespace)):
                        cells = [
                            clean_text("".join(node.text or "" for node in cell.findall(".//w:t", namespace)))
                            for cell in row.findall("./w:tc", namespace)
                        ]
                        if any(cells):
                            document.elements.append(SourceElement(
                                id=self._element_id(document, "row", f"openxml:{name}:{table_index}:{row_index}"),
                                kind=ElementKind.TABLE_ROW,
                                text=" | ".join(cells),
                                row_number=row_index + 1,
                                metadata={"table_id": table_id, "cells": cells, "parser": "openxml"},
                            ))
        self._append_text_chunks(
            document,
            "\n".join(paragraphs),
            location="openxml-page:1",
            page_number=1,
            metadata={"parser": "openxml", "image_extraction": False},
        )

    def _parse_pptx(
        self, path: Path, document: ParsedDocument, artifact_dir: Path, *, fast: bool = False
    ) -> None:
        presentation = Presentation(path)
        image_dir = artifact_dir / "figures"
        image_dir.mkdir(exist_ok=True)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape_index, shape in enumerate(slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    value = clean_text(shape.text)
                    if value:
                        texts.append(value)
                if getattr(shape, "has_table", False):
                    table_id = self._element_id(document, "table", f"slide:{slide_index}:{shape_index}")
                    document.elements.append(
                        SourceElement(
                            id=table_id,
                            kind=ElementKind.TABLE,
                            slide_number=slide_index,
                            text=f"Slide {slide_index} table",
                        )
                    )
                    for row_index, row in enumerate(shape.table.rows):
                        cells = [clean_text(cell.text) for cell in row.cells]
                        document.elements.append(
                            SourceElement(
                                id=self._element_id(
                                    document, "row", f"slide:{slide_index}:{shape_index}:{row_index}"
                                ),
                                kind=ElementKind.TABLE_ROW,
                                slide_number=slide_index,
                                row_number=row_index + 1,
                                text=" | ".join(cells),
                                metadata={"table_id": table_id, "cells": cells},
                            )
                        )
                if not fast and getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    content = shape.image.blob
                    if len(content) < 8000:
                        continue
                    digest = hashlib.sha256(content).hexdigest()
                    extension = shape.image.ext or "png"
                    target = image_dir / f"{digest}.{extension}"
                    target.write_bytes(content)
                    document.elements.append(
                        SourceElement(
                            id=self._element_id(document, "figure", f"slide:{slide_index}:{digest}"),
                            kind=ElementKind.FIGURE,
                            slide_number=slide_index,
                            image_path=str(target),
                            metadata={"sha256": digest},
                        )
                    )
            try:
                notes = clean_text(slide.notes_slide.notes_text_frame.text)
                if notes:
                    texts.append(f"Speaker notes: {notes}")
            except Exception:
                pass
            self._append_text_chunks(
                document,
                "\n".join(texts),
                location=f"slide:{slide_index}",
                slide_number=slide_index,
                page_number=slide_index,
            )

    def _parse_workbook(self, path: Path, document: ParsedDocument) -> None:
        book = pd.ExcelFile(path)
        for sheet_index, sheet_name in enumerate(book.sheet_names):
            frame = pd.read_excel(path, sheet_name=sheet_name, header=None)
            table_id = self._element_id(document, "table", f"sheet:{sheet_index}:{sheet_name}")
            document.elements.append(
                SourceElement(
                    id=table_id,
                    kind=ElementKind.TABLE,
                    text=f"Workbook sheet: {sheet_name}",
                    sheet_name=sheet_name,
                    metadata={"rows": int(frame.shape[0]), "columns": int(frame.shape[1])},
                )
            )
            for row_index, row in frame.iterrows():
                cells = []
                for value in row.tolist():
                    if pd.isna(value):
                        cells.append("")
                    elif hasattr(value, "isoformat"):
                        cells.append(value.isoformat())
                    else:
                        cells.append(clean_text(str(value)))
                if not any(cells):
                    continue
                document.elements.append(
                    SourceElement(
                        id=self._element_id(document, "row", f"sheet:{sheet_index}:{row_index}"),
                        kind=ElementKind.TABLE_ROW,
                        text=" | ".join(cells)[:20000],
                        sheet_name=sheet_name,
                        row_number=int(row_index) + 1,
                        metadata={"table_id": table_id, "cells": cells},
                    )
                )

    def _infer_title(self, document: ParsedDocument) -> str:
        for element in document.elements:
            if element.kind == ElementKind.TEXT and element.text:
                candidates = [clean_text(line) for line in element.text.splitlines() if clean_text(line)]
                candidates = [line for line in candidates if 5 <= len(line) <= 300]
                if candidates:
                    return candidates[0]
        return Path(document.file_name).stem

    def _initial_publications(self, document: ParsedDocument) -> list[PublicationBoundary]:
        pages = [element.page_number for element in document.elements if element.page_number]
        max_page = max(pages, default=1)
        is_container = bool(max_page >= 50 and (
            document.category in {"Журналы", "Материалы конференций"}
            or re.search(r"proceedings|journal|журнал|конферен", document.file_name, re.I)
        ))
        return [
            PublicationBoundary(
                title=document.title or document.file_name,
                start_page=1,
                end_page=max_page,
                confidence=0.25 if is_container else 0.95,
                needs_review=is_container,
            )
        ]


def validate_manifest(manifest_path: str | Path, corpus_root: str | Path | None = None) -> list[dict]:
    manifest_path = Path(manifest_path)
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    root = Path(corpus_root or data.get("corpus_root", "/corpus"))
    documents = data.get("documents", [])
    if len(documents) != 25:
        raise ValueError(f"Pilot manifest must contain exactly 25 documents, got {len(documents)}")
    counts: dict[str, int] = {}
    resolved = []
    for item in documents:
        category = item["category"]
        counts[category] = counts.get(category, 0) + 1
        relative = Path(item["path"])
        if relative.is_absolute() or ".." in relative.parts:
            raise ValueError(f"Unsafe manifest path: {relative}")
        path = root / relative
        if not path.is_file():
            raise FileNotFoundError(path)
        if path.suffix.lower() not in DocumentParser.SUPPORTED:
            raise ValueError(f"Unsupported manifest source: {path}")
        resolved.append({**item, "resolved_path": str(path)})
    if set(counts.values()) != {5} or len(counts) != 5:
        raise ValueError(f"Manifest must contain five documents in each of five categories: {counts}")
    return resolved
